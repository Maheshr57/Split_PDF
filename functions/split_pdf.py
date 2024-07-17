from transformers import BartForConditionalGeneration, BartTokenizer
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from fpdf import FPDF

def read_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def read_text_from_pdf(file_path, page_ranges):
    doc = fitz.open(file_path)
    text = ""
    
    last_page = 0
    for start, end in page_ranges:
        for page_num in range(start, end):
            text += doc.load_page(page_num).get_text()
        last_page = end

    # Automatically add remaining pages to the last split
    if last_page < doc.page_count:
        for page_num in range(last_page, doc.page_count):
            text += doc.load_page(page_num).get_text()
    
    return text

def read_text_from_word(file_path):
    doc = docx.Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text

def read_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def summarize_text_chunk(text_chunk, max_length=150, min_length=30):
    tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')
    model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')

    inputs = tokenizer([text_chunk], max_length=1024, return_tensors='pt', truncation=True)
    summary_ids = model.generate(inputs['input_ids'], num_beams=4, max_length=max_length, min_length=min_length, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)

    return summary

def split_text_into_chunks(text, chunk_size=1024):
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

def summarize_large_text(text, chunk_size=1024, max_length=150, min_length=30):
    chunks = split_text_into_chunks(text, chunk_size)
    summaries = [summarize_text_chunk(chunk, max_length, min_length) for chunk in chunks]
    return ' '.join(summaries)

def summarize_file(file_path, page_ranges):
    if file_path.endswith('.txt'):
        text = read_text_from_txt(file_path)
    elif file_path.endswith('.pdf'):
        text = read_text_from_pdf(file_path, page_ranges)
    elif file_path.endswith('.docx'):
        text = read_text_from_word(file_path)
    elif file_path.endswith('.pptx'):
        text = read_text_from_ppt(file_path)
    else:
        raise ValueError("Unsupported file format")
    
    return summarize_large_text(text)

def save_summary_to_pdf(summary, output_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    # Split the summary into lines to fit the PDF width
    lines = summary.split('\n')
    for line in lines:
        pdf.multi_cell(0, 10, line.encode('latin-1', 'replace').decode('latin-1'))

    pdf.output(output_path)

# Example usage
def main():
    file_path = 'path/to/your/file.pdf'  # Replace with your file path
    page_ranges = [(0, 5), (5, 10)]  # Example page ranges, automatically handles remaining pages

    # Summarize file
    summary = summarize_file(file_path, page_ranges)
    print("Summary:\n", summary)

    # Save summary to PDF
    output_pdf_path = 'path/to/save/summary.pdf'  # Path to save the summary PDF
    save_summary_to_pdf(summary, output_pdf_path)
    print(f"Summary saved to {output_pdf_path}")

if __name__ == "__main__":
    main()
